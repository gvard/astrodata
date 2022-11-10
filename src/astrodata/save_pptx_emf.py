"""Script for collecting images, converting SVG to EMF format and
adding resulting files to pptx presentation using python-pptx package.
It is similar to Photo album generating in PowerPoint, but automatically.
Needs installed Inkscape.
"""

import subprocess
import os
import glob
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches
from PIL import Image


Image.MAX_IMAGE_PIXELS = None # avoiding PIL.Image.DecompressionBombError
INKSCAPE_PTH = "c:/Program Files/Inkscape/bin/inkscape.exe" # edit it in Linux
today = datetime.now()
MONTH, YEAR = today.strftime("%m"), today.year

plots_dir = os.path.join(os.pardir, os.pardir, 'plots')
files_to_add = []

for fn in glob.iglob(plots_dir + "../../plots/*/*.png"):
    files_to_add.append(fn)
for fn in glob.iglob(plots_dir + "../../plots/*/*.svg"):
    emf_fname = fn.rstrip("svg") + "emf"
    print(f"Converting {fn} to EMF format...")
    subprocess.run([INKSCAPE_PTH, fn, '--export-filename', emf_fname], check=True)
    files_to_add.append(emf_fname)

print(f"Adding {len(files_to_add)} images to pptx presentation...")

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
wdth, hght = Inches(13.333), Inches(7.5)
left = top = Inches(0)
prs.slide_width = wdth
prs.slide_height = hght

for img_filename in files_to_add:
    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(img_filename, left, top, wdth, hght)

prs.save(f'astrostats-{YEAR}-{MONTH}.pptx')
